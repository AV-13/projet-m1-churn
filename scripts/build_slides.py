"""Génère le support de présentation (.pptx) éditable du projet.

Identité visuelle « Signal » (cohérente avec le dashboard). Usage :
    python scripts/build_slides.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "figures"
OUT = ROOT / "docs" / "04-presentation.pptx"

INK = RGBColor(0x1A, 0x22, 0x33)
MUTED = RGBColor(0x6B, 0x73, 0x85)
PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xE1, 0x1D, 0x48)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF3, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333
n = [0]  # compteur de slides


def _set(p, size, color, bold=False, font="Calibri"):
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font


def new_slide():
    return prs.slides.add_slide(BLANK)


def footer(slide):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = "Prédiction du churn · Sujet 2 · EFREI"
    _set(p, 9, MUTED)
    nb = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(0.8), Inches(0.35))
    pp = nb.text_frame.paragraphs[0]
    pp.text = str(n[0])
    pp.alignment = PP_ALIGN.RIGHT
    _set(pp, 9, MUTED)


def title_bar(slide, title, eyebrow=None):
    if eyebrow:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        p.text = eyebrow.upper()
        p.alignment = PP_ALIGN.CENTER
        _set(p, 11, PRIMARY, bold=True)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(12.1), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    _set(p, 28, INK, bold=True)
    # Trait d'accent centré sous le titre (cohérent quel que soit le logiciel).
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.07), Inches(1.62),
                                  Inches(1.2), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = PRIMARY; line.line.fill.background()


def bullets(slide, items, left=0.7, top=1.95, width=12.0, size=18, gap=10):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.8))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        txt, *rest = it if isinstance(it, tuple) else (it,)
        lvl = rest[0] if rest else 0
        color = rest[1] if len(rest) > 1 else INK
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("    " * lvl) + ("– " if lvl else "•  ") + txt
        _set(p, size - (2 if lvl else 0), color, bold=(lvl == 0 and len(rest) > 1))
        p.space_after = Pt(gap)


def image_right(slide, path, left=7.2, top=1.95, width=5.6):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def content(title, items, eyebrow=None, image=None, text_width=12.0, size=18):
    n[0] += 1
    s = new_slide()
    title_bar(s, title, eyebrow)
    bullets(s, items, width=text_width, size=size)
    if image:
        image_right(s, image)
    footer(s)
    return s


# --------------------------------------------------------------------------- #
# 1. Slide de titre
# --------------------------------------------------------------------------- #
n[0] += 1
s = new_slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
band.fill.solid(); band.fill.fore_color.rgb = PRIMARY; band.line.fill.background()
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.5), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Prédiction du churn client"; _set(p, 44, INK, bold=True)
p = tf.add_paragraph(); p.text = "Système intelligent de rétention et d'évaluation du risque de revenus"
_set(p, 19, MUTED); p.space_before = Pt(6)
p = tf.add_paragraph(); p.text = "Sujet 2 · M1 Dev. Manager Full Stack · Data Science · EFREI 2025-26"
_set(p, 13, PRIMARY, bold=True); p.space_before = Pt(18)
p = tf.add_paragraph(); p.text = "[Noms du groupe] · 26 juin 2026"
_set(p, 13, MUTED); p.space_before = Pt(4)

# --------------------------------------------------------------------------- #
# 2 → 11. Contenu
# --------------------------------------------------------------------------- #
content("Le problème : la fuite des clients", [
    "Les entreprises par abonnement (SaaS, télécom…) perdent des clients chaque mois (churn).",
    "Retenir un client coûte bien moins cher que d'en acquérir un nouveau.",
    "Objectif : prédire le risque de départ, l'expliquer, et le rendre actionnable.",
    "Utilisateur cible : responsable marketing / CRM.",
], eyebrow="Contexte & objectif")

content("Un jeu de données fortement déséquilibré", [
    "customer_churn.csv — 10 000 clients, 30 variables (numériques + catégorielles).",
    "Cible : churn (0 = reste / 1 = part).",
    ("Seulement 10,2 % de churners → ratio de déséquilibre 8,8:1.", 1, RED),
    ("L'accuracy est trompeuse → on privilégie Recall, F1 et PR-AUC.", 1, INK),
    "Signal métier fort : ≥ 2 paiements ratés → churn de 21 à 33 % (vs 8,8 %).",
], eyebrow="Données & EDA")

content("Une démarche rigoureuse, pas seulement un score", [
    "Pipeline de préparation anti-fuite (appris sur le train uniquement).",
    "Validation croisée stratifiée (5 plis) ; métrique de sélection : PR-AUC.",
    "Gestion du déséquilibre : SMOTE, class_weight, et ajustement du seuil de décision.",
    "4 modèles comparés, du plus simple au plus complexe :",
    ("Régression logistique → Random Forest → Gradient Boosting → MLP (Deep Learning).", 1),
], eyebrow="Méthodologie")

# Slide résultats avec tableau
n[0] += 1
s = new_slide()
title_bar(s, "Comparaison des modèles", "Résultats")
rows = [("Modèle", "PR-AUC", "ROC-AUC"),
        ("Random Forest  ★", "0,267", "0,794"),
        ("HistGradientBoosting", "0,259", "0,789"),
        ("Régression logistique", "0,226", "0,715"),
        ("MLP (Deep Learning)", "0,167", "0,631")]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.7), Inches(2.0),
                         Inches(7.2), Inches(3.0)).table
tbl.columns[0].width = Inches(4.0); tbl.columns[1].width = Inches(1.6)
tbl.columns[2].width = Inches(1.6)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
            _set(para, 14, WHITE, bold=True)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r % 2 == 0 else WHITE
            _set(para, 13, INK, bold=(r == 1))
side = s.shapes.add_textbox(Inches(8.3), Inches(2.0), Inches(4.4), Inches(4.0))
tf = side.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "À retenir"; _set(p, 15, PRIMARY, bold=True)
p = tf.add_paragraph()
p.text = "Le MLP (Deep Learning) est le moins performant."
_set(p, 15, RED, bold=True); p.space_before = Pt(8)
p = tf.add_paragraph()
p.text = ("Sur des données tabulaires, les modèles à arbres surpassent le réseau "
          "de neurones — pour un coût bien moindre.")
_set(p, 13, INK); p.space_before = Pt(6)
p = tf.add_paragraph()
p.text = "→ Le Deep Learning n'est pas toujours supérieur."
_set(p, 13, MUTED, bold=True); p.space_before = Pt(6)
footer(s)

content("Le piège du déséquilibre : le seuil de décision", [
    ("Au seuil par défaut 0,5 : recall = 0 — aucun churner détecté ! (malgré 90 % d'accuracy)", 0, RED),
    ("Au seuil optimisé 0,19 : recall = 78 %, F1 = 0,37.", 0, GREEN),
    "Choix orienté métier : un faux négatif (client perdu) coûte plus cher",
    ("qu'un faux positif (une simple sollicitation marketing).", 1),
    "→ Régler le seuil est aussi déterminant que choisir le modèle.",
], eyebrow="Résultats", image=FIG / "confusion_matrix.png", text_width=6.4)

content("Pourquoi un client part", [
    "Facteurs dominants identifiés par le modèle :",
    ("Satisfaction (CSAT) — facteur n°1, de loin.", 1),
    ("Échecs de paiement — le risque triple dès 2 échecs.", 1),
    ("Ancienneté et engagement (logins, usage).", 1),
    "SHAP explique chaque prédiction individuelle → décisions justifiables.",
], eyebrow="Interprétabilité", image=FIG / "permutation_importance.png", text_width=6.4)

content("Démonstration : le dashboard décisionnel", [
    "Jauge de risque claire + verdict (faible / modéré / élevé).",
    "Signaux du profil lus en un coup d'œil.",
    "Simulation d'une action de rétention (impact sur le risque).",
    "Le dashboard consomme l'API (architecture réaliste).",
], eyebrow="Démo", image=FIG / "dashboard_scored.png", text_width=5.6)

content("Une solution industrialisée et reproductible", [
    "API REST (FastAPI) : /predict, /health, /model-info — validation + gestion d'erreurs.",
    "Dashboard (Streamlit) qui consomme l'API : architecture Front → API → Modèle.",
    "Pipeline de préprocessing sérialisé avec le modèle (cohérence train/prod).",
    "Reproductibilité : Makefile, requirements, tests pytest, dépôt Git versionné.",
], eyebrow="Architecture")

content("Limites et recommandations", [
    ("Limites", 0, PRIMARY),
    ("Pouvoir prédictif modéré (données synthétiques) ; précision 24 % au seuil retenu ; dérive temporelle.", 1),
    ("Recommandations métier", 0, PRIMARY),
    ("Cibler les clients à ≥ 2 paiements ratés ; surveiller la satisfaction (CSAT/NPS) ; sécuriser les clients récents.", 1),
    ("Recommandations techniques", 0, PRIMARY),
    ("Feature engineering ; suivi de dérive et réentraînement périodique.", 1),
], eyebrow="Recul critique", size=16)

content("Conclusion", [
    "D'un dataset brut à une solution complète : explicable, industrialisée, exploitable.",
    "Deux enseignements clés :",
    ("Sur données tabulaires, un modèle simple (Random Forest) surpasse le Deep Learning.", 1),
    ("En contexte déséquilibré, le réglage du seuil est aussi déterminant que le modèle.", 1),
    "Merci — questions ?",
], eyebrow="Synthèse")

prs.save(OUT)
print("Présentation générée :", OUT, f"({n[0]} slides)")
